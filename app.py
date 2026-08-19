# -*- coding: utf-8 -*-
"""
Clarilex（旧Minivibe）
--------------------------------------------------------------
どんな長さ・テーマの原稿でも、自由な枚数のスライドへ組み立てられる
Streamlitアプリ。AIによる自動生成、長文の一括流し込み、手動での
追加・編集・並べ替えを組み合わせて使える。

各スライドは「タイトル」「メインテキスト（長文対応）」「サブテキスト」
「画像（アップロード／URL／キーワード検索）」の4要素で構成される。
メインテキストは文字数に応じて自動でフォントサイズが縮小され、
それでも収まりきらない場合はプレビュー内でスクロールする。

起動方法:
    pip install -r requirements.txt
    streamlit run app.py

必要なAPIキー（すべて任意。手動編集・長文流し込みだけならAPIキー不要）:
    - Cohere / Google Gemini / Groq のいずれか（AIによるスライド自動生成用）
        Cohere : https://dashboard.cohere.com/
        Gemini : https://aistudio.google.com/
        Groq   : https://console.groq.com/
    - Unsplash Access Key（キーワード検索モードの画像精度向上用）
        https://unsplash.com/developers
"""

import io
import re
import json
import uuid
import html
import base64
from urllib.parse import quote

import streamlit as st

# ページ設定は最初のStreamlitコマンドである必要がある
st.set_page_config(page_title="Clarilex", page_icon="", layout="wide")

import pandas as pd

# python-pptx / Pillow はアプリの核となる必須ライブラリのため、
# 未インストールの場合はここで分かりやすいエラーを出して停止する。
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from PIL import Image as PILImage
except ImportError:
    st.error(
        "`python-pptx` または `Pillow` がインストールされていません。\n\n"
        "ターミナルで `pip install python-pptx Pillow` を実行してから、アプリを再起動してください。"
    )
    st.stop()


# ============================================================
# 定数・テーマ定義
# ============================================================

THEMES = {
    "Apple Black": {"bg": "#000000", "text": "#FFFFFF", "sub": "#98989D", "metric": "#FFFFFF", "border": "#1D1D1F"},
    "Apple White": {"bg": "#FFFFFF", "text": "#1D1D1F", "sub": "#86868B", "metric": "#1D1D1F", "border": "#E5E5E7"},
    "Cyber":       {"bg": "#05050A", "text": "#EAFBFF", "sub": "#7DFFEA", "metric": "#00FFC6", "border": "#1B1B2F"},
    "Minimal Grey":{"bg": "#F5F5F7", "text": "#1D1D1F", "sub": "#6E6E73", "metric": "#1D1D1F", "border": "#D2D2D7"},
}

FONT_STACK_CSS = (
    "-apple-system, BlinkMacSystemFont, 'SF Pro Display', 'Hiragino Sans', "
    "'Yu Gothic', 'Segoe UI', Roboto, sans-serif"
)
PPTX_FONT_NAME = "Yu Gothic"  # 日本語表示に強いフォント（未インストール環境ではOS側で代替されます）

DEFAULT_COHERE_MODEL = "command-a-03-2025"
DEFAULT_GEMINI_MODEL = "gemini-3.6-flash"
DEFAULT_GROQ_MODEL = "openai/gpt-oss-120b"

UNSPLASH_SEARCH_URL = "https://api.unsplash.com/search/photos"
DEFAULT_IMAGE_CANDIDATE_COUNT = 5
MAX_UPLOAD_BYTES = 8 * 1024 * 1024  # 8MB

IMAGE_SOURCE_LABELS = {
    "upload": "A：ファイルからアップロード",
    "url": "B：画像URLを指定",
    "keyword": "C：キーワード検索",
}
IMAGE_SOURCE_KEYS = ["upload", "url", "keyword"]

# メインテキストの文字数に応じたフォント縮小の段階（文字数, 倍率）
MAIN_TEXT_SCALE_TIERS = [
    (8, 1.00), (20, 0.88), (40, 0.72), (70, 0.58), (110, 0.46),
    (170, 0.36), (260, 0.28), (400, 0.22), (600, 0.18),
]
MAIN_TEXT_SCALE_FALLBACK = 0.15


# ============================================================
# ユーティリティ
# ============================================================

def hexstr(h: str) -> str:
    """'#FFFFFF' -> 'FFFFFF' に変換（RGBColor.from_string用）"""
    return h.lstrip("#").upper()


def main_text_scale(text: str) -> float:
    """メインテキストの文字数に応じたフォント縮小係数(0.0〜1.0)を返す。
    短い一文は大きく力強く、長い文章は自動的に読みやすいサイズへ縮小する。"""
    length = len(str(text or "").strip())
    for max_len, scale in MAIN_TEXT_SCALE_TIERS:
        if length <= max_len:
            return scale
    return MAIN_TEXT_SCALE_FALLBACK


def new_slide_dict(title: str = "", main_text: str = "", subtext: str = "",
                    image_keyword: str = "", show_image: bool = True,
                    image_source: str = "keyword") -> dict:
    return {
        "id": str(uuid.uuid4()),
        "title": title,
        "main_text": main_text,
        "subtext": subtext,
        "show_image": show_image,
        "image_source": image_source,        # "upload" | "url" | "keyword"
        "image_keyword": image_keyword,       # C: キーワード検索
        "image_keyword_locked": False,
        "image_url": "",
        "image_credit": "",
        "image_manual_url": "",               # B: URL直接指定
        "image_upload_data_uri": "",          # A: アップロード
        "image_upload_name": "",
    }


def clear_slide_image(slide: dict):
    """スライドから画像本体を取り除く（検索キーワード自体・選択中モードは残す）"""
    slide["image_url"] = ""
    slide["image_credit"] = ""
    slide["image_manual_url"] = ""
    slide["image_upload_data_uri"] = ""
    slide["image_upload_name"] = ""


def image_file_to_data_uri(uploaded_file) -> str:
    """st.file_uploaderで受け取ったファイルをBase64データURIに変換する"""
    raw_bytes = uploaded_file.getvalue()
    if len(raw_bytes) > MAX_UPLOAD_BYTES:
        raise RuntimeError(f"画像サイズが大きすぎます（{MAX_UPLOAD_BYTES // (1024*1024)}MB以下にしてください）。")
    mime = uploaded_file.type or "image/png"
    b64 = base64.b64encode(raw_bytes).decode("utf-8")
    return f"data:{mime};base64,{b64}"


def decode_data_uri_to_bytes(data_uri: str) -> bytes:
    """Base64データURIから画像バイナリを取り出す（pptx埋め込み用）"""
    try:
        _, encoded = data_uri.split(",", 1)
        return base64.b64decode(encoded)
    except Exception as e:
        raise RuntimeError(f"アップロード画像のデコードに失敗しました: {e}")


def get_slide_image_src(slide: dict):
    """現在選択中のモードに応じた (画像src, 種類, クレジット) を返す。
    種類は 'data_uri' | 'url'。画像が無ければ (None, None, None)。"""
    source = slide.get("image_source", "keyword")

    if source == "upload":
        data_uri = slide.get("image_upload_data_uri") or ""
        return (data_uri, "data_uri", None) if data_uri else (None, None, None)

    if source == "url":
        url = (slide.get("image_manual_url") or "").strip()
        return (url, "url", None) if url else (None, None, None)

    url = slide.get("image_url") or ""
    return (url, "url", slide.get("image_credit")) if url else (None, None, None)


# ============================================================
# 長文の自動分割（AI不要・「✂️ 長文を貼り付けて一括スライド化」用）
# ============================================================

def split_long_text(text: str, mode: str, chunk_size: int = 300):
    """長文を複数の (タイトル, 本文) チャンクへ分割する。
    mode: 'blank_line'（空行区切り） | 'heading'（見出し区切り） | 'char_count'（文字数で均等分割）"""
    text = (text or "").strip()
    if not text:
        return []

    if mode == "heading":
        parts = re.split(r"(?m)^(?=#{1,6}\s|■|【.+】|\d+[\.\)]\s)", text)
        raw_chunks = [p for p in parts if p.strip()]
    elif mode == "char_count":
        chunk_size = max(50, int(chunk_size))
        raw_chunks = [text[i:i + chunk_size] for i in range(0, len(text), chunk_size)]
    else:  # blank_line
        raw_chunks = re.split(r"\n\s*\n+", text)

    raw_chunks = [c.strip() for c in raw_chunks if c.strip()]
    if not raw_chunks:
        raw_chunks = [text]

    result = []
    for i, chunk in enumerate(raw_chunks, start=1):
        lines = chunk.splitlines()
        first_line = lines[0].strip() if lines else ""
        if first_line and len(first_line) <= 30 and re.match(r"^(#{1,6}\s|■|【|\d+[\.\)]\s)", first_line):
            title = re.sub(r"^(#{1,6}\s|■|【)", "", first_line).strip("】 ")
            body = "\n".join(lines[1:]).strip() or first_line
        elif first_line and len(first_line) <= 24 and len(lines) > 1:
            title = first_line
            body = "\n".join(lines[1:]).strip()
        else:
            title = f"スライド {i}"
            body = chunk
        result.append((title, body))
    return result


# ============================================================
# AIプロンプト生成
# ============================================================

def build_prompt(source_text: str, n_slides: int, image_lang: str = "English (推奨)") -> str:
    instruction = (
        "あなたは優秀なエンジニアです。冗長な説明を排除し、"
        "AppleのKeynoteプレゼンテーションのような最小限かつ洗練された表現に変換すること。"
    )

    if image_lang.startswith("English"):
        image_lang_note = "英語で2〜5単語程度。Unsplash等の海外ストックフォトサービスでの検索精度を優先すること。"
    else:
        image_lang_note = "日本語で2〜5単語程度の具体的なキーワード。"

    return f"""{instruction}

以下の原稿を、プレゼンテーション用に約 {n_slides} 枚のスライドへ分解してください。
各スライドは「1枚につき1テーマ」を意識しつつ、次の4要素で構成してください。

1. title         : スライドの見出し（日本語で15〜20文字程度の短い一文）
2. main_text     : そのスライドの本文。短い強調フレーズでも、数文程度の説明文でも構いません。
                    原稿の内容を薄めすぎず、必要な情報は残しつつ簡潔にまとめること。
3. subtext       : 補足説明（1〜2行程度。無ければ空文字でも可）
4. image_keyword : スライドの文脈・テーマを象徴する、具体的で視覚的な画像検索キーワード。
                   {image_lang_note}
                   抽象的な概念語ではなく、実際に写真として存在しうる具体物・情景を指定すること。

【出力形式（厳守）】
説明文・前置き・Markdownのコードブロック(```)は一切使わず、次のJSON配列のみを出力してください。

[
  {{"title": "...", "main_text": "...", "subtext": "...", "image_keyword": "..."}},
  {{"title": "...", "main_text": "...", "subtext": "...", "image_keyword": "..."}}
]

【変換対象の原稿】
{source_text}
"""


def extract_json_array(raw_text: str):
    """AIの応答からJSON配列部分を頑健に抽出する"""
    if not raw_text or not raw_text.strip():
        raise ValueError("AIからの応答が空でした。")

    text = raw_text.strip()
    text = re.sub(r"^```[a-zA-Z]*\n?", "", text)
    text = re.sub(r"```$", "", text).strip()

    start = text.find("[")
    end = text.rfind("]")
    if start == -1 or end == -1 or end < start:
        raise ValueError(f"応答からJSON配列が見つかりませんでした。応答冒頭: {text[:200]}")

    json_str = text[start:end + 1]
    try:
        data = json.loads(json_str)
    except json.JSONDecodeError as e:
        raise ValueError(f"JSONの解析に失敗しました: {e}\n応答冒頭: {json_str[:300]}")

    if not isinstance(data, list):
        raise ValueError("AIの応答がリスト（配列）形式ではありませんでした。")
    return data


# ============================================================
# AI呼び出し（Cohereをメイン、Gemini / Groqを補助として実装）
# ============================================================

def call_cohere(prompt: str, api_key: str, model: str = DEFAULT_COHERE_MODEL) -> str:
    try:
        import cohere
    except ImportError:
        raise RuntimeError("`cohere` ライブラリが未インストールです。`pip install cohere` を実行してください。")

    last_err = None
    try:
        co = cohere.ClientV2(api_key=api_key)
        resp = co.chat(model=model, messages=[{"role": "user", "content": prompt}])
        content = getattr(getattr(resp, "message", None), "content", None)
        if isinstance(content, list) and len(content) > 0 and hasattr(content[0], "text"):
            return content[0].text
        if content:
            return str(content)
        return str(resp)
    except Exception as e:
        last_err = e

    try:
        co = cohere.Client(api_key)
        resp = co.chat(model=model, message=prompt)
        if hasattr(resp, "text"):
            return resp.text
        return str(resp)
    except Exception as e2:
        raise RuntimeError(f"Cohere APIの呼び出しに失敗しました。(ClientV2: {last_err} / Client: {e2})")


def call_gemini(prompt: str, api_key: str, model: str = DEFAULT_GEMINI_MODEL) -> str:
    try:
        import google.generativeai as genai
    except ImportError:
        raise RuntimeError("`google-generativeai` ライブラリが未インストールです。`pip install google-generativeai` を実行してください。")

    try:
        genai.configure(api_key=api_key)
        gm = genai.GenerativeModel(model)
        resp = gm.generate_content(prompt)
        text = getattr(resp, "text", None)
        if text:
            return text
        if getattr(resp, "candidates", None):
            parts = resp.candidates[0].content.parts
            return "".join(getattr(p, "text", "") for p in parts)
        raise RuntimeError("Geminiから有効な応答が得られませんでした。")
    except Exception as e:
        raise RuntimeError(f"Gemini APIの呼び出しに失敗しました: {e}")


def call_groq(prompt: str, api_key: str, model: str = DEFAULT_GROQ_MODEL) -> str:
    try:
        from groq import Groq
    except ImportError:
        raise RuntimeError("`groq` ライブラリが未インストールです。`pip install groq` を実行してください。")

    try:
        client = Groq(api_key=api_key)
        resp = client.chat.completions.create(
            model=model,
            messages=[
                {"role": "system", "content": (
                    "あなたは優秀なエンジニアです。冗長な説明を排除し、"
                    "AppleのKeynoteプレゼンテーションのような最小限かつ洗練された表現に変換すること。"
                )},
                {"role": "user", "content": prompt},
            ],
            temperature=0.7,
        )
        return resp.choices[0].message.content
    except Exception as e:
        raise RuntimeError(f"Groq APIの呼び出しに失敗しました: {e}")


# ============================================================
# 画像検索・選定（Cモード：キーワード検索）
# ============================================================

def _fetch_unsplash_candidates(keyword: str, access_key: str, count: int = DEFAULT_IMAGE_CANDIDATE_COUNT,
                                timeout: int = 10):
    try:
        import requests
    except ImportError:
        raise RuntimeError("`requests` ライブラリが未インストールです。`pip install requests` を実行してください。")

    resp = requests.get(
        UNSPLASH_SEARCH_URL,
        params={"query": keyword, "per_page": count, "orientation": "landscape"},
        headers={"Authorization": f"Client-ID {access_key}"},
        timeout=timeout,
    )
    resp.raise_for_status()
    results = (resp.json() or {}).get("results") or []
    if not results:
        raise RuntimeError(f"Unsplashで '{keyword}' に一致する画像が見つかりませんでした。")

    candidates = []
    for photo in results:
        url = (photo.get("urls") or {}).get("regular") or (photo.get("urls") or {}).get("small")
        if not url:
            continue
        photographer = (photo.get("user") or {}).get("name", "Unsplash")
        description = photo.get("description") or photo.get("alt_description") or ""
        candidates.append({"url": url, "credit": f"Photo by {photographer} on Unsplash", "description": description.strip()})
    if not candidates:
        raise RuntimeError("Unsplashの応答に有効な画像URLが含まれていませんでした。")
    return candidates


def select_best_image_with_cohere(slide: dict, candidates: list, cohere_key: str,
                                   cohere_model: str = DEFAULT_COHERE_MODEL) -> dict:
    """候補画像の中から、スライドの文脈に最もふさわしい1枚をCohereに選ばせる。"""
    if not candidates:
        raise ValueError("候補画像がありません。")
    if len(candidates) == 1:
        return candidates[0]
    if not cohere_key:
        raise ValueError("Cohere APIキーが未設定のため画像選定できません。")

    lines = [f"{i}. {c.get('description') or '(説明文なし)'}" for i, c in enumerate(candidates, start=1)]
    main_text_preview = (slide.get("main_text", "") or "")[:200]

    prompt = f"""あなたはプレゼンテーションの画像選定担当です。
以下のスライド内容に、視覚的に最もふさわしい画像を候補の中から1つだけ選んでください。

【スライド内容】
タイトル: {slide.get('title', '')}
本文: {main_text_preview}
補足: {slide.get('subtext', '')}

【画像候補】
{chr(10).join(lines)}

最も適した候補の番号を、半角数字1つだけで答えてください。説明・記号・文章は一切不要です。
"""
    raw = call_cohere(prompt, cohere_key, cohere_model)
    match = re.search(r"\d+", raw or "")
    if not match:
        raise ValueError(f"Cohereの応答から番号を抽出できませんでした: {raw!r}")

    idx = int(match.group())
    if not (1 <= idx <= len(candidates)):
        idx = 1
    return candidates[idx - 1]


def fetch_image_for_slide(slide: dict, unsplash_access_key: str = None, cohere_key: str = None,
                           cohere_model: str = DEFAULT_COHERE_MODEL,
                           candidate_count: int = DEFAULT_IMAGE_CANDIDATE_COUNT,
                           keyword_locked: bool = False):
    """キーワード検索モード用に画像を選び、(url, credit) を返す。
    固定OFF: Unsplash複数候補からCohereが選定。固定ON: 先頭候補を直接採用（精度優先）。
    Unsplash未設定/失敗時はダミープレースホルダー画像にフォールバックする。"""
    keyword = (slide.get("image_keyword") or "").strip()
    if not keyword:
        raise ValueError("画像検索キーワードが空です。")

    if unsplash_access_key:
        try:
            fetch_count = 1 if keyword_locked else candidate_count
            candidates = _fetch_unsplash_candidates(keyword, unsplash_access_key, count=fetch_count)
            if keyword_locked or len(candidates) == 1:
                chosen = candidates[0]
            else:
                try:
                    chosen = select_best_image_with_cohere(slide, candidates, cohere_key, cohere_model)
                except Exception:
                    chosen = candidates[0]
            return chosen["url"], chosen["credit"]
        except Exception:
            pass

    seed = quote(keyword)
    return f"https://picsum.photos/seed/{seed}/900/600", None


@st.cache_data(show_spinner=False, ttl=3600)
def download_image_bytes(url: str) -> bytes:
    """pptx埋め込み用に画像バイナリを取得する。同一URLはセッション内でキャッシュされる。"""
    try:
        import requests
    except ImportError:
        raise RuntimeError("`requests` ライブラリが未インストールです。`pip install requests` を実行してください。")
    resp = requests.get(url, timeout=15)
    resp.raise_for_status()
    return resp.content


def generate_slides_with_ai(source_text, engine, cohere_key, gemini_key, groq_key,
                             cohere_model, gemini_model, groq_model, n_slides,
                             image_lang="English (推奨)", auto_fetch_images=True,
                             unsplash_key=None, image_candidate_count=DEFAULT_IMAGE_CANDIDATE_COUNT):
    """AIでスライド本文を生成し、必要なら画像も自動取得する。
    戻り値: (slides: list[dict], image_warnings: list[str])"""
    prompt = build_prompt(source_text, n_slides, image_lang)

    if engine == "Cohere":
        if not cohere_key:
            raise ValueError("Cohere APIキーが入力されていません。サイドバーに入力してください。")
        raw = call_cohere(prompt, cohere_key, cohere_model)
    elif engine == "Google Gemini":
        if not gemini_key:
            raise ValueError("Gemini APIキーが入力されていません。サイドバーに入力してください。")
        raw = call_gemini(prompt, gemini_key, gemini_model)
    elif engine == "Groq":
        if not groq_key:
            raise ValueError("Groq APIキーが入力されていません。サイドバーに入力してください。")
        raw = call_groq(prompt, groq_key, groq_model)
    else:
        raise ValueError(f"不明なエンジンです: {engine}")

    st.session_state["_last_ai_raw"] = raw

    data = extract_json_array(raw)
    slides = []
    for item in data:
        if not isinstance(item, dict):
            continue
        slides.append(new_slide_dict(
            title=str(item.get("title", "")).strip(),
            main_text=str(item.get("main_text", "")).strip(),
            subtext=str(item.get("subtext", "")).strip(),
            image_keyword=str(item.get("image_keyword", "")).strip(),
            image_source="keyword",
        ))

    if not slides:
        raise ValueError("AIの応答からスライドを1件も抽出できませんでした。")

    image_warnings = []
    if auto_fetch_images:
        for s in slides:
            if not s.get("image_keyword", "").strip():
                continue
            try:
                url, credit = fetch_image_for_slide(
                    s, unsplash_key, cohere_key, cohere_model, image_candidate_count, keyword_locked=False,
                )
                s["image_url"] = url
                s["image_credit"] = credit
            except Exception as e:
                image_warnings.append(f"「{(s['title'] or s['image_keyword'])[:15]}」: {e}")

    return slides, image_warnings


# ============================================================
# スライド操作（st.session_state管理）
# ============================================================

def move_slide(index: int, direction: int):
    slides = st.session_state.slides
    new_index = index + direction
    if 0 <= new_index < len(slides):
        slides[index], slides[new_index] = slides[new_index], slides[index]


def delete_slide(index: int):
    slides = st.session_state.slides
    if 0 <= index < len(slides):
        slides.pop(index)


def add_blank_slide():
    st.session_state.slides.append(new_slide_dict(title="新しいスライド", main_text="", subtext=""))


def clear_all_images():
    for s in st.session_state.slides:
        clear_slide_image(s)


# ============================================================
# プレビューカード（HTML/CSS）
#   ・長文でもはみ出さないよう、文字数に応じた自動縮小 + スクロール安全弁を実装
# ============================================================

def _render_text_block_html(slide: dict, theme: dict, compact: bool) -> str:
    """タイトル・メインテキスト・サブテキストのHTMLブロックを生成する。
    compact=Trueは画像ありの2カラムレイアウト時（やや小さめのサイズにする）。"""
    title = html.escape(slide.get("title", "") or "")
    main_text = html.escape(slide.get("main_text", "") or "")
    subtext = html.escape(slide.get("subtext", "") or "")

    base_main_px = 52 if compact else 68
    min_main_px = 15
    scale = main_text_scale(slide.get("main_text", ""))
    main_px = max(min_main_px, int(base_main_px * scale))

    return f"""
    <div style="flex:1; min-width:0; text-align:left;">
        <div style="
            font-size:{18 if compact else 20}px; font-weight:700; letter-spacing:0.3px;
            color:{theme['text']}; margin-bottom:12px; opacity:0.92;
        ">{title}</div>
        <div style="
            font-size:{main_px}px; font-weight:800; line-height:1.3;
            color:{theme['metric']}; margin:4px 0 14px 0;
            white-space:pre-wrap; word-break:break-word; overflow-wrap:anywhere;
            max-height:280px; overflow-y:auto; padding-right:6px;
        ">{main_text}</div>
        <div style="
            font-size:13px; font-weight:400; color:{theme['sub']}; line-height:1.5;
            white-space:pre-wrap; word-break:break-word; overflow-wrap:anywhere;
            max-height:100px; overflow-y:auto;
        ">{subtext}</div>
    </div>
    """


def render_preview_card_html(slide: dict, theme: dict) -> str:
    src, _kind, credit = get_slide_image_src(slide)
    has_image = bool(slide.get("show_image")) and bool(src)

    text_block = _render_text_block_html(slide, theme, compact=has_image)

    if has_image:
        image_src_safe = html.escape(src)
        credit_html = (
            f'<div style="font-size:10px; color:{theme["sub"]}; opacity:0.75; margin-top:6px; text-align:right;">{html.escape(credit)}</div>'
            if credit else ""
        )
        image_block = f"""
        <div style="flex:1; min-width:0;">
            <img src="{image_src_safe}" style="width:100%; height:220px; object-fit:cover; border-radius:16px; display:block;" />
            {credit_html}
        </div>
        """
        inner = text_block + image_block
        layout_style = "display:flex; align-items:flex-start; gap:28px;"
    else:
        inner = text_block
        layout_style = "display:block;"

    return f"""
    <div style="
        background:{theme['bg']};
        color:{theme['text']};
        border:1px solid {theme['border']};
        border-radius:24px;
        padding:36px;
        margin-bottom:18px;
        box-shadow:0 10px 30px rgba(0,0,0,0.18);
        font-family:{FONT_STACK_CSS};
        {layout_style}
    ">
        {inner}
    </div>
    """


# ============================================================
# PowerPoint (.pptx) 生成
# ============================================================

def add_pptx_textbox(slide, text: str, left_in: float, top_in: float, width_in: float, height_in: float,
                      size_pt: int, bold: bool, color_rgb, align=PP_ALIGN.LEFT):
    """テキストボックスを追加する共通ヘルパー。改行(\\n)は段落として反映する。"""
    box = slide.shapes.add_textbox(Inches(left_in), Inches(top_in), Inches(width_in), Inches(height_in))
    tf = box.text_frame
    tf.word_wrap = True

    lines = (text or "").split("\n") or [""]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size_pt)
        run.font.bold = bold
        run.font.name = PPTX_FONT_NAME
        run.font.color.rgb = color_rgb
    return box


def add_cover_image(slide, image_bytes: bytes, left_in: float, top_in: float, width_in: float, height_in: float):
    """CSSのobject-fit:coverのように、アスペクト比を保ったまま指定枠いっぱいに
    画像をクロップ配置する（python-pptxのcrop_*プロパティを利用）。"""
    img_w = img_h = None
    try:
        with PILImage.open(io.BytesIO(image_bytes)) as im:
            img_w, img_h = im.size
    except Exception:
        pass

    pic = slide.shapes.add_picture(io.BytesIO(image_bytes), Inches(left_in), Inches(top_in),
                                    width=Inches(width_in), height=Inches(height_in))

    if img_w and img_h:
        target_ratio = width_in / height_in
        img_ratio = img_w / img_h
        if img_ratio > target_ratio:
            crop = (1 - target_ratio / img_ratio) / 2
            pic.crop_left, pic.crop_right = crop, crop
        elif img_ratio < target_ratio:
            crop = (1 - img_ratio / target_ratio) / 2
            pic.crop_top, pic.crop_bottom = crop, crop
    return pic


def build_pptx(slides: list, theme_name: str):
    """編集後のスライドデータから16:9のPowerPointを生成する。
    戻り値: (buffer: io.BytesIO, warnings: list[str])
    ※ pptxは静的なスライドのためスクロールできない。極端に長いmain_textは
      自動縮小はされるが、それでも収まりきらない場合は見切れうる点に注意。"""
    theme = THEMES[theme_name]
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    bg_color = RGBColor.from_string(hexstr(theme["bg"]))
    text_color = RGBColor.from_string(hexstr(theme["text"]))
    sub_color = RGBColor.from_string(hexstr(theme["sub"]))
    metric_color = RGBColor.from_string(hexstr(theme["metric"]))

    warnings = []

    for slide_data in slides:
        slide = prs.slides.add_slide(blank_layout)
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = bg_color

        title_text = slide_data.get("title", "")
        short_label = (title_text or slide_data.get("image_keyword") or "")[:15]

        img_src, img_kind, img_credit = get_slide_image_src(slide_data)
        has_image = bool(slide_data.get("show_image")) and bool(img_src)
        image_bytes = None
        if has_image:
            try:
                image_bytes = decode_data_uri_to_bytes(img_src) if img_kind == "data_uri" else download_image_bytes(img_src)
            except Exception as e:
                warnings.append(f"「{short_label}」の画像取得に失敗しました: {e}")
                has_image = False

        scale = main_text_scale(slide_data.get("main_text", ""))

        if has_image:
            # ---- 左：テキスト / 右：画像 の2カラムレイアウト ----
            text_left, text_width = 0.6, 6.0
            img_left, img_top, img_width, img_height = 6.9, 0.6, 5.83, 6.3

            add_pptx_textbox(slide, title_text, text_left, 0.7, text_width, 0.7, 24, True, text_color)
            add_pptx_textbox(slide, slide_data.get("main_text", ""), text_left, 1.5, text_width, 4.3,
                              max(int(40 * scale), 12), True, metric_color)
            add_pptx_textbox(slide, slide_data.get("subtext", ""), text_left, 5.9, text_width, 1.0, 14, False, sub_color)

            try:
                add_cover_image(slide, image_bytes, img_left, img_top, img_width, img_height)
            except Exception as e:
                warnings.append(f"「{short_label}」の画像配置に失敗しました: {e}")

            if img_credit:
                add_pptx_textbox(slide, img_credit, img_left, img_top + img_height + 0.05, img_width, 0.3,
                                  9, False, sub_color, align=PP_ALIGN.RIGHT)
        else:
            # ---- 画像なし：中央寄せの単一カラムレイアウト ----
            col_left, col_width = 2.0, 9.33
            add_pptx_textbox(slide, title_text, col_left, 0.8, col_width, 0.8, 30, True, text_color)
            add_pptx_textbox(slide, slide_data.get("main_text", ""), col_left, 1.7, col_width, 4.2,
                              max(int(56 * scale), 14), True, metric_color)
            add_pptx_textbox(slide, slide_data.get("subtext", ""), col_left, 6.1, col_width, 1.0, 15, False, sub_color)

    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer, warnings


# ============================================================
# セッション状態の初期化
# ============================================================

if "slides" not in st.session_state:
    st.session_state.slides = [
        new_slide_dict(
            title="shorter. faster. bolder.",
            main_text="10x",
            subtext="短いフレーズは、自動的に大きく力強く表示されます。",
            image_keyword="minimalist product design",
        ),
        new_slide_dict(
            title="長文でも崩れません",
            main_text=(
                "メインテキストが長くなると、フォントサイズが自動的に縮小されます。"
                "それでも収まりきらない場合は、カード内でスクロールして読めるため、"
                "レポートや記事の本文をそのまま貼り付けても、レイアウトが崩れる心配はありません。"
            ),
            subtext="左の「編集」パネルから自由に書き換えられます。",
            image_source="keyword",
        ),
    ]


# ============================================================
# サイドバー UI
# ============================================================

with st.sidebar:
    st.header("⚙️ 設定")

    st.subheader("🧩 クイック操作")
    if st.button("➕ 新しいスライドを追加", use_container_width=True, key="sidebar_add_slide"):
        add_blank_slide()
        st.rerun()

    st.subheader("🔑 APIキー（任意）")
    st.caption("APIキーはこのセッション内でのみ使用され、サーバーに保存されません。手動編集や長文の一括流し込みだけならAPIキーは不要です。")

    cohere_api_key = st.text_input("Cohere APIキー（メイン・推奨）", type="password", key="cohere_api_key")
    st.caption("スライド自動生成に加え、キーワード検索モード（固定OFF時）の画像選定にも使用されます。")

    with st.expander("補助エンジン（Google Gemini / Groq）"):
        gemini_api_key = st.text_input("Google Gemini APIキー", type="password", key="gemini_api_key")
        groq_api_key = st.text_input("Groq APIキー", type="password", key="groq_api_key")

    st.subheader("🖼️ 画像検索（Cモード用・任意）")
    unsplash_api_key = st.text_input(
        "Unsplash Access Key", type="password", key="unsplash_api_key",
        help="https://unsplash.com/developers で無料取得できます。",
    )
    st.caption("未入力の場合は、キーワードから決定的に選ばれるダミープレースホルダー画像（picsum.photos）が使われます。")

    engine = st.selectbox("🤖 生成に使うAIエンジン", ["Cohere", "Google Gemini", "Groq"], index=0)

    st.subheader("🎨 デザインテーマ")
    theme_name = st.selectbox("テーマを選択", list(THEMES.keys()), index=0, label_visibility="collapsed")

    with st.expander("詳細設定（モデル名 / 画像 / デバッグ）"):
        cohere_model = st.text_input("Cohereモデル", value=DEFAULT_COHERE_MODEL)
        gemini_model = st.text_input("Geminiモデル", value=DEFAULT_GEMINI_MODEL)
        groq_model = st.text_input("Groqモデル", value=DEFAULT_GROQ_MODEL)
        image_lang = st.selectbox("画像検索キーワードの言語", ["English (推奨)", "日本語"], index=0)
        image_candidate_count = st.slider("画像候補数（Cohereによる選定用・固定OFF時のみ）", 2, 8, DEFAULT_IMAGE_CANDIDATE_COUNT)
        show_debug = st.checkbox("AIの生応答を表示する", value=False)


# ============================================================
# メイン UI
# ============================================================

st.title("Clarilex")
st.caption("どんな原稿でも、自由な枚数のスライドへ。")

tab_ai, tab_paste = st.tabs(["AIで自動生成", "長文を貼り付けて一括スライド化（AI不要）"])

with tab_ai:
    raw_text = st.text_area(
        "変換したい原稿・長文を貼り付けてください",
        height=180,
        placeholder="ここに元原稿を入力...（例：ハンムラビ法典から近代法の成立までを説明した長文など）",
        key="ai_raw_text",
    )
    gen_col1, gen_col2 = st.columns([3, 1])
    with gen_col2:
        n_slides = st.number_input("目安枚数", min_value=1, max_value=40, value=5, step=1)
    with gen_col1:
        st.write("")
        generate_clicked = st.button("AIでスライド生成（既存スライドを置き換え）", type="primary", use_container_width=True)

    auto_fetch_images = st.checkbox("生成と同時に画像も自動取得する（キーワード検索モード）", value=True)
    st.caption("40枚を超える非常に長い原稿は、右の「長文を貼り付けて一括スライド化」タブの利用がおすすめです（枚数上限なし）。")

    if generate_clicked:
        if not raw_text or not raw_text.strip():
            st.warning("原稿を入力してください。")
        else:
            with st.spinner("AIが原稿を削ぎ落とし、画像を選定しています..."):
                try:
                    new_slides, img_warnings = generate_slides_with_ai(
                        raw_text, engine, cohere_api_key, gemini_api_key, groq_api_key,
                        cohere_model, gemini_model, groq_model, n_slides,
                        image_lang, auto_fetch_images, unsplash_api_key, image_candidate_count,
                    )
                    st.session_state.slides = new_slides
                    st.success(f"{len(new_slides)}枚のスライドを生成しました。")
                    if img_warnings:
                        st.warning("一部のスライドで画像取得に失敗しました：\n" + "\n".join(img_warnings))
                    st.rerun()
                except Exception as e:
                    st.error(f"生成に失敗しました: {e}")

    if show_debug and st.session_state.get("_last_ai_raw"):
        with st.expander("AIの生応答（デバッグ）"):
            st.code(st.session_state["_last_ai_raw"])

with tab_paste:
    st.caption("APIキー不要。長文をそのまま貼り付けると、指定した方法で自動的に複数スライドへ分割します。枚数に上限はありません。")
    bulk_text = st.text_area(
        "ここに長文（レポート・記事など）を貼り付け",
        height=200,
        placeholder="長文をそのまま貼り付けてください。空行や見出しごとに自動でスライド分割できます。",
        key="bulk_paste_text",
    )
    split_mode = st.radio(
        "分割方法",
        options=["blank_line", "heading", "char_count"],
        format_func=lambda k: {
            "blank_line": "空行で段落ごとに分割",
            "heading": "見出し（#・■・【】・番号）で分割",
            "char_count": "文字数で均等分割",
        }[k],
        horizontal=True,
        key="split_mode",
    )
    if split_mode == "char_count":
        chunk_size = st.number_input("1スライドあたりの目安文字数", min_value=50, max_value=2000, value=300, step=50)
    else:
        chunk_size = 300

    paste_col1, paste_col2 = st.columns([2, 1])
    with paste_col1:
        append_mode = st.checkbox("既存のスライドの後ろに追加する（OFFですべて置き換え）", value=True)
    with paste_col2:
        split_clicked = st.button("スライドに流し込む", type="primary", use_container_width=True)

    if split_clicked:
        if not bulk_text or not bulk_text.strip():
            st.warning("テキストを入力してください。")
        else:
            chunks = split_long_text(bulk_text, split_mode, chunk_size)
            new_slides = [new_slide_dict(title=t, main_text=b) for t, b in chunks]
            if not new_slides:
                st.warning("スライドを作成できませんでした。テキストを確認してください。")
            else:
                if append_mode:
                    st.session_state.slides.extend(new_slides)
                else:
                    st.session_state.slides = new_slides
                st.success(f"{len(new_slides)}枚のスライドを作成しました。")
                st.rerun()

st.divider()

left_col, right_col = st.columns([1, 1])

# ---------------- 左カラム：編集エリア（アコーディオン） ----------------
with left_col:
    st.subheader("編集")

    top_btn_col1, top_btn_col2 = st.columns(2)
    with top_btn_col1:
        if st.button("➕新しいスライドを追加", use_container_width=True, key="main_add_slide"):
            add_blank_slide()
            st.rerun()
    with top_btn_col2:
        if st.button("すべての画像を削除", use_container_width=True):
            clear_all_images()
            st.rerun()

    if not st.session_state.slides:
        st.info("まだスライドがありません。AIで生成する、長文を貼り付ける、または上のボタンで追加してください。")
    else:
        for i, slide in enumerate(st.session_state.slides):
            sid = slide["id"]
            label = slide["title"].strip() or "(未入力)"
            with st.expander(f"スライド {i + 1}：{label[:20]}", expanded=(i == 0)):
                slide["title"] = st.text_input("スライドタイトル", value=slide["title"], key=f"title_{sid}")
                slide["main_text"] = st.text_area(
                    "メインテキスト（長文対応）", value=slide["main_text"], key=f"main_{sid}", height=160,
                    help="短い一言でも長い文章でも、プレビュー側で自動的にフォントサイズが調整されます。",
                )
                slide["subtext"] = st.text_area("サブテキスト（補足）", value=slide["subtext"], key=f"subtext_{sid}", height=70)

                st.markdown("画像")
                src_col, show_col = st.columns([3, 1])
                with src_col:
                    slide["image_source"] = st.radio(
                        "画像の指定方法", options=IMAGE_SOURCE_KEYS,
                        index=IMAGE_SOURCE_KEYS.index(slide.get("image_source", "keyword")),
                        format_func=lambda k: IMAGE_SOURCE_LABELS[k],
                        key=f"imgsrc_{sid}", horizontal=True, label_visibility="collapsed",
                    )
                with show_col:
                    slide["show_image"] = st.checkbox("表示する", value=slide.get("show_image", True), key=f"showimg_{sid}")

                if slide["image_source"] == "upload":
                    uploaded = st.file_uploader(
                        "画像ファイルを選択", type=["png", "jpg", "jpeg", "gif", "webp", "bmp"],
                        key=f"upload_{sid}", label_visibility="collapsed",
                    )
                    if uploaded is not None:
                        try:
                            slide["image_upload_data_uri"] = image_file_to_data_uri(uploaded)
                            slide["image_upload_name"] = uploaded.name
                        except Exception as e:
                            st.error(f"画像の読み込みに失敗しました: {e}")

                elif slide["image_source"] == "url":
                    slide["image_manual_url"] = st.text_input(
                        "画像URL", value=slide.get("image_manual_url", ""),
                        key=f"manualurl_{sid}", label_visibility="collapsed",
                        placeholder="https://example.com/photo.jpg",
                    )

                else:  # keyword
                    kw_col1, kw_col2 = st.columns([2, 1])
                    with kw_col1:
                        slide["image_keyword"] = st.text_input(
                            "画像検索キーワード", value=slide.get("image_keyword", ""),
                            key=f"imgkw_{sid}", label_visibility="collapsed",
                            placeholder="例: ancient stone law tablet",
                        )
                    with kw_col2:
                        slide["image_keyword_locked"] = st.checkbox(
                            "固定", value=slide.get("image_keyword_locked", False), key=f"lockkw_{sid}",
                            help="ONにすると、Cohereによる候補選定をスキップし、このキーワードの検索結果を直接使用します（精度・再現性優先）。",
                        )
                    if st.button("🔄 画像を検索/更新", key=f"fetchimg_{sid}", use_container_width=True):
                        kw = (slide["image_keyword"] or "").strip()
                        if not kw:
                            st.warning("画像検索キーワードを入力してください。")
                        else:
                            spinner_msg = "指定キーワードで検索中..." if slide["image_keyword_locked"] else "Cohereが画像を選定中..."
                            with st.spinner(spinner_msg):
                                try:
                                    url, credit = fetch_image_for_slide(
                                        slide, unsplash_api_key, cohere_api_key, cohere_model,
                                        image_candidate_count, keyword_locked=slide["image_keyword_locked"],
                                    )
                                    slide["image_url"] = url
                                    slide["image_credit"] = credit
                                    st.rerun()
                                except Exception as e:
                                    st.error(f"画像の取得に失敗しました: {e}")

                preview_src, preview_kind, preview_credit = get_slide_image_src(slide)
                if preview_src:
                    if preview_kind == "data_uri":
                        try:
                            st.image(decode_data_uri_to_bytes(preview_src), caption=slide.get("image_upload_name") or None, use_container_width=True)
                        except Exception as e:
                            st.error(f"プレビューの表示に失敗しました: {e}")
                    else:
                        st.image(preview_src, caption=preview_credit or None, use_container_width=True)

                    if st.button("この画像を削除", key=f"delimg_{sid}", use_container_width=True):
                        clear_slide_image(slide)
                        st.rerun()

                b1, b2, b3 = st.columns(3)
                if b1.button("⬆️ 上へ", key=f"up_{sid}", use_container_width=True):
                    move_slide(i, -1)
                    st.rerun()
                if b2.button("⬇️ 下へ", key=f"down_{sid}", use_container_width=True):
                    move_slide(i, 1)
                    st.rerun()
                if b3.button("🗑️ スライドを削除", key=f"del_{sid}", use_container_width=True):
                    delete_slide(i)
                    st.rerun()

        with st.expander("全体構成一覧（テーブル表示）"):
            df = pd.DataFrame(st.session_state.slides)[["title", "main_text", "subtext", "image_source", "image_keyword"]].copy()
            df["main_text"] = df["main_text"].apply(lambda t: (t[:40] + "…") if isinstance(t, str) and len(t) > 40 else t)
            df.index = df.index + 1
            df.columns = ["タイトル", "メインテキスト", "補足", "画像モード", "画像キーワード"]
            st.dataframe(df, use_container_width=True)

# ---------------- 右カラム：プレビューエリア ----------------
with right_col:
    st.subheader("プレビュー")

    theme = THEMES[theme_name]

    if not st.session_state.slides:
        st.info("プレビューするスライドがありません。")
    else:
        st.caption(f"全 {len(st.session_state.slides)} 枚")
        for slide in st.session_state.slides:
            st.markdown(render_preview_card_html(slide, theme), unsafe_allow_html=True)

        try:
            pptx_buffer, pptx_warnings = build_pptx(st.session_state.slides, theme_name)
            if pptx_warnings:
                st.warning("PowerPoint生成時に一部の画像を取得できませんでした：\n" + "\n".join(pptx_warnings))
            st.download_button(
                "⬇️ PowerPointをダウンロード（.pptx）",
                data=pptx_buffer,
                file_name="clarilex_presentation.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                use_container_width=True,
                type="primary",
            )
        except Exception as e:
            st.error(f"PowerPointの生成に失敗しました: {e}")
